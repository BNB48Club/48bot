# -*- coding: utf-8 -*-
import sys
import convertapi
from datetime import datetime
from pptx import Presentation

def genPNG(title,content):
    template = "BNB48snapshot.pptx"

    prs = Presentation(template)
    shape=prs.slides[0].shapes
    for i in range(0,len(shape)):
        if(shape[i].has_text_frame):
            for paragraph in shape[i].text_frame.paragraphs:
                for run in paragraph.runs:
                    if 'title' in run.text:
                        run.text = title
                    if 'content' in run.text:
                        run.text = content
                    if 'date' in run.text:
                        run.text = datetime.now().strftime("%m/%d/%Y")

    prs.save("temp.pptx")
    convertapi.api_secret = 'Ea4zdB42huH3eRMX'
    convertapi.convert('png', {
        'File': 'temp.pptx'
    }, from_format = 'pptx').save_files('temp.png')
    return "temp.png"
